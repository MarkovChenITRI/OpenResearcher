"""
Enhanced file content extraction utilities
Support for PDF, DOCX, PPTX, and images
"""
import os
from pathlib import Path

def extract_pdf_content(filepath):
    """Extract text from PDF"""
    try:
        import PyPDF2
        with open(filepath, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = []
            for page in reader.pages:
                text.append(page.extract_text())
            return '\n\n'.join(text)
    except ImportError:
        return "[PDF] PyPDF2 not installed. Run: pip install PyPDF2"
    except Exception as e:
        return f"[PDF] Error: {e}"


def extract_docx_content(filepath):
    """Extract text from DOCX"""
    try:
        import docx
        doc = docx.Document(filepath)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        return '\n\n'.join(text)
    except ImportError:
        return "[DOCX] python-docx not installed. Run: pip install python-docx"
    except Exception as e:
        return f"[DOCX] Error: {e}"


def extract_pptx_content(filepath):
    """Extract text from PPTX"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        text = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Slide {slide_num} ==="]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if len(slide_text) > 1:
                text.append('\n'.join(slide_text))
        return '\n\n'.join(text)
    except ImportError:
        return "[PPTX] python-pptx not installed. Run: pip install python-pptx"
    except Exception as e:
        return f"[PPTX] Error: {e}"


def extract_image_content(filepath):
    """Extract text from image using OCR (optional)"""
    try:
        # Optional: Use OCR
        # from PIL import Image
        # import pytesseract
        # image = Image.open(filepath)
        # text = pytesseract.image_to_string(image)
        # return text
        
        # For now, just return filename
        return f"[Image: {Path(filepath).name}]\n(OCR not implemented. Install pytesseract for OCR support)"
    except Exception as e:
        return f"[Image] Error: {e}"


def extract_file_content(filepath):
    """Main dispatcher for file content extraction"""
    ext = Path(filepath).suffix.lower()
    
    extractors = {
        '.md': lambda f: open(f, 'r', encoding='utf-8').read(),
        '.txt': lambda f: open(f, 'r', encoding='utf-8').read(),
        '.pdf': extract_pdf_content,
        '.docx': extract_docx_content,
        '.pptx': extract_pptx_content,
        '.png': extract_image_content,
        '.jpg': extract_image_content,
        '.jpeg': extract_image_content,
    }
    
    extractor = extractors.get(ext)
    if extractor:
        try:
            return extractor(filepath)
        except Exception as e:
            return f"Error extracting {ext}: {e}"
    else:
        return f"Unsupported file type: {ext}"


# Test if executed directly
if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        filepath = sys.argv[1]
        print(f"Extracting content from: {filepath}")
        print("=" * 60)
        content = extract_file_content(filepath)
        print(content)
        print("=" * 60)
        print(f"Extracted {len(content)} characters")
    else:
        print("Usage: python file_utils.py <filepath>")
